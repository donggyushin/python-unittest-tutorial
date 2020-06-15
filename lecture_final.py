from PyPDF2 import PdfFileReader, PdfFileWriter, PdfFileMerger
import xlsxwriter as xw
import openpyxl as xr
from fpdf import FPDF
import csv
import os
from docx import Document
from pptx import Presentation
# csv ㅍㅏ일 읽기


def read():
    f = open('ex.csv', 'rt')
    reader = csv.DictReader(f)

    data = list(reader)
    f.close()

# csv 파일 쓰기


def write():
    f = open('ex.csv', 'wt')

    writer = csv.writer(f, delimiter='\t', lineterminator='\n')
    writer.writerow(('sr', 'names', 'grades'))

    f.close()


# 엑셀파일 읽기


def readExcel(filePath):
    # 수식이 계산된 후의 값을 원하기 때문에 data_only=True 기본값은 False임
    workbook = xr.load_workbook(filePath, data_only=True)
    workbook['sheetName']

# 엑셀파일 작성


def writeExcel(filePath):

    book = xw.Workbook(filePath)
    # name은 있어도 되고 없어도 됨.
    sheet = book.add_worksheet(name="sheetName")
    # 0번째 행에 첫번째 컬럼에 sdf 문자를 작성
    row = 0
    col = 1
    sheet.write(row, col, 'sdf')


# pdf 파일 읽기


def readPdf():
    pdf = open('ex.pdf', 'rb')
    reader = PdfFileReader(pdf)
    # pdf 파일 정보 얻기
    reader.getDocumentInfo()


def writePdf():
    # format의 기본값 letter
    pdf = FPDF()
    pdf.add_page()
    # 높이, 너비, 텍스트, 라인 간격, 정렬
    pdf.cell(200, 100, txt="Hello world", ln=1, align='C')
    pdf.output('ex.pdf')
    pdf.close()

# 파일 병합


def mergePdf():
    merger = PdfFileMerger()
    files = [x for x in os.listdir() if x.endswith('pdf')]

    for fname in sorted(files):
        merger.append(PdfFileReader(open(fname, 'rb')))

    f1 = open('ex1.pdf', 'rb')
    f2 = open('ex2.pdf', 'rb')
    merger.append(PdfFileReader(f1))
    merger.append(PdfFileReader(f2))

    merger.write('ex.pdf')
    merger.close()


# 워드 문서 읽기
def readDocx():
    doc = Document('ex.docx')
    doc.core_properties

# 워드 문서 작성


def writeDocx():
    doc = Document()
    r = doc.add_paragraph('A plain paragraph having some')

    # 밑에 r 아닌가?
    p.add_run('bold words').bold = True

    doc.save('ex.docx')


def readPPT():
    p = Presentation('ex.pptx')


def writePPT():
    p = Presentation()

    slide = p.slides.add_slide(p.slide_layouts[3])
    slide.shapes.title.text = 'Hi'
    slide.placeholders[1].text = 'Python'

    p.save('ex.pptx')
